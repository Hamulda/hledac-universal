"""
RustMiner - Lightweight Data Mining for M1 Optimization
======================================================
Memory-efficient content extraction using lightweight libraries.

Optimization Strategy:
- Uses trafilex (Rust-based) or traflatura in minimal mode
- No large DOM trees - streaming processing
- Minimal memory footprint for M1 8GB
"""

import logging
import re
from dataclasses import field
from operator import itemgetter
from typing import Any
from urllib.parse import urljoin, urlparse

from hledac.universal.compat.msgspec_gc_compat import Struct

# selectolax — strict import with fallback
try:
    from selectolax.parser import HTMLParser

    SELECTOLAX_AVAILABLE = True
except ImportError:
    HTMLParser = None
    SELECTOLAX_AVAILABLE = False

# lxml — strict import with fallback
try:
    from lxml import html as lxml_html

    LXML_AVAILABLE = True
except ImportError:
    lxml_html = None
    LXML_AVAILABLE = False

# nh3 — Rust HTML sanitizer, 9× faster than regex, M1 8GB friendly
try:
    import nh3

    NH3_AVAILABLE = True
except ImportError:
    nh3 = None
    NH3_AVAILABLE = False

logger = logging.getLogger(__name__)

# html_text_fast — strict import with fallback
try:
    from hledac.universal.utils.html_text_fast import html_to_text_fast

    _CANONICAL_HTML_TEXT_AVAILABLE = True
except ImportError:
    try:
        from hledac.universal.utils.html_text_fast import html_to_text_fast

        _CANONICAL_HTML_TEXT_AVAILABLE = True
    except ImportError:
        html_to_text_fast = None
        _CANONICAL_HTML_TEXT_AVAILABLE = False
# nh3-clean fallback patterns — used only when nh3 is unavailable
_NH3_FALLBACK_PATTERNS: list[tuple[re.Pattern, str]] = [
    (re.compile("<script[^>]*>.*?</script>", flags=re.DOTALL | re.IGNORECASE), " "),
    (re.compile("<style[^>]*>.*?</style>", flags=re.DOTALL | re.IGNORECASE), " "),
    (re.compile("<head[^>]*>.*?</head>", flags=re.DOTALL | re.IGNORECASE), " "),
    (re.compile("<nav[^>]*>.*?</nav>", flags=re.DOTALL | re.IGNORECASE), " "),
    (re.compile("<footer[^>]*>.*?</footer>", flags=re.DOTALL | re.IGNORECASE), " "),
    (re.compile("<header[^>]*>.*?</header>", flags=re.DOTALL | re.IGNORECASE), " "),
    (re.compile("<aside[^>]*>.*?</aside>", flags=re.DOTALL | re.IGNORECASE), " "),
    (re.compile("<[^>]+>"), " "),
    (re.compile("&nbsp;"), " "),
    (re.compile("&amp;"), "&"),
    (re.compile("&lt;"), "<"),
    (re.compile("&gt;"), ">"),
    (re.compile("&quot;"), '"'),
    (re.compile("&apos;"), "'"),
    (re.compile(r"\s+"), " "),
]
_TAG_STRIP_PATTERN = re.compile(r"<[^>]+>")
_MULTI_SPACE_PATTERN = re.compile(r"\s+")


class MiningResult(Struct):
    """Result of content mining operation"""

    content: str
    title: str = ""
    url: str = ""
    metadata: dict[str, Any] = None
    success: bool = True
    error: str | None = None

    def __post_init__(self) -> None:
        if self.metadata is None:
            self.metadata = {}


class RustMiner:
    """
    Lightweight content miner using Rust-backed libraries.

    Strategy:
    1. Try trafilex (Rust, fastest) - minimal DOM
    2. Fallback to traflatura (minimal mode) - streaming
    3. Ultimate fallback to regex (no dependencies)
    """

    __slots__ = ("_trafilex_available", "_traflatura_available", "_use_ultimate_fallback", "prefer_rust")

    def __init__(self, prefer_rust: bool = True) -> None:
        """
        Initialize RustMiner.

        Args:
            prefer_rust: Prefer Rust-based libraries (trafilex) over Python
        """
        self.prefer_rust = prefer_rust
        self._trafilex_available = self._check_trafilex()
        self._traflatura_available = self._check_traflatura()
        self._use_ultimate_fallback = not (self._trafilex_available or self._traflatura_available)
        logger.info("✓ RustMiner initialized")
        logger.info(f"  trafilex: {self._trafilex_available}")
        logger.info(f"  traflatura: {self._traflatura_available}")
        logger.info(f"  fallback: {self._use_ultimate_fallback}")

    def _check_trafilex(self) -> bool:
        """Check if trafilex (Rust-based) is available"""
        try:
            import trafilex

            return True
        except ImportError:
            return False

    def _check_traflatura(self) -> bool:
        """Check if traflatura is available"""
        try:
            import traflatura

            return True
        except ImportError:
            return False

    def mine_html(self, html_content: str, url: str = "", include_metadata: bool = False) -> MiningResult:
        """
        Extract clean text from HTML with minimal memory usage.

        Args:
            html_content: Raw HTML content
            url: Source URL for metadata
            include_metadata: Extract metadata (title, author, date)

        Returns:
            MiningResult with extracted content
        """
        try:
            if not html_content or not isinstance(html_content, str):
                return MiningResult(content="", url=url, success=False, error="Invalid HTML content")
            # OSINT-03: Bound input size before parsing to prevent OOM on M1 8GB.
            if len(html_content) > self._MAX_HTML_INPUT_SIZE:
                html_content = html_content[: self._MAX_HTML_INPUT_SIZE]
            logger.info("[MINER] Extraction started...")
            if self.prefer_rust and self._trafilex_available:
                return self._mine_with_trafilex(html_content, url, include_metadata)
            elif self._traflatura_available:
                return self._mine_with_traflatura(html_content, url, include_metadata)
            else:
                return self._mine_with_fallback(html_content, url, include_metadata)
        except Exception as e:
            logger.error(f"HTML mining failed: {e}")
            return MiningResult(content="", url=url, success=False, error=str(e))

    def _mine_with_trafilex(self, html_content: str, url: str, include_metadata: bool) -> MiningResult:
        """
        Mine using trafilex (Rust-based, minimal memory).

        trafilex is a Rust wrapper that processes HTML in streaming mode
        without building large DOM trees.
        """
        try:
            import trafilex

            result = trafilex.extract(html_content, include_comments=False, include_tables=False, no_fallback=False)
            title = result.title or ""
            content = result.content or ""
            if not content:
                content = self._clean_html_basic(html_content)
            metadata = {}
            if include_metadata:
                metadata.update(
                    {"title": title, "source": url, "method": "trafilex (Rust)", "char_count": len(content)}
                )
                jsonld = self.extract_jsonld(html_content)
                if jsonld:
                    metadata["jsonld"] = jsonld
            logger.debug(f"trafilex extracted {len(content)} chars from {url}")
            return MiningResult(content=content, title=title, url=url, metadata=metadata, success=True)
        except Exception as e:
            logger.warning(f"trafilex failed: {e}, falling back")
            if self._traflatura_available:
                return self._mine_with_traflatura(html_content, url, include_metadata)
            return self._mine_with_fallback(html_content, url, include_metadata)

    def _mine_with_traflatura(self, html_content: str, url: str, include_metadata: bool) -> MiningResult:
        """
        Mine using traflatura in minimal mode.

        Memory optimization:
        - disable_comments: Don't store comment nodes
        - no_tables: Skip table extraction (expensive)
        - include_tables: False to save memory
        - deduplicate: True to reduce memory
        """
        try:
            import traflatura

            settings = traflatura.settings.use_settings()
            settings.output_format = "txt"
            settings.include_comments = False
            settings.include_tables = False
            settings.include_formatting = False
            settings.include_images = False
            settings.deduplicate = True
            result = traflatura.extract(
                html_content,
                settings=settings,
                url=url,
                include_comments=False,
                include_tables=False,
                no_fallback=False,
            )
            title = self._extract_title_fallback(html_content)
            content = result or ""
            if not content:
                content = self._clean_html_basic(html_content)
            metadata = {}
            if include_metadata:
                metadata.update(
                    {"title": title, "source": url, "method": "traflatura (minimal)", "char_count": len(content)}
                )
                jsonld = self.extract_jsonld(html_content)
                if jsonld:
                    metadata["jsonld"] = jsonld
            logger.debug(f"traflatura extracted {len(content)} chars from {url}")
            return MiningResult(content=content, title=title, url=url, metadata=metadata, success=True)
        except Exception as e:
            logger.warning(f"traflatura failed: {e}, using fallback")
            return self._mine_with_fallback(html_content, url, include_metadata)

    def _mine_with_fallback(self, html_content: str, url: str, include_metadata: bool) -> MiningResult:
        """
        Ultimate fallback using regex-based extraction.

        No dependencies - pure Python regex for maximum compatibility.
        """
        try:
            title = self._extract_title_fallback(html_content)
            content = self._clean_html_basic(html_content)
            metadata = {}
            if include_metadata:
                metadata.update({"title": title, "source": url, "method": "regex fallback", "char_count": len(content)})
                jsonld = self.extract_jsonld(html_content)
                if jsonld:
                    metadata["jsonld"] = jsonld
            logger.debug(f"Regex fallback extracted {len(content)} chars from {url}")
            return MiningResult(content=content, title=title, url=url, metadata=metadata, success=True)
        except Exception as e:
            logger.error(f"Fallback extraction failed: {e}")
            return MiningResult(content="", url=url, success=False, error=str(e))

    def _clean_html_basic(self, html: str) -> str:
        """
        Basic HTML cleaning — nh3 primary + tag stripping (Rust sanitizer, 9× faster).

        Cascade: nh3 (strips dangerous content) + tag strip → equivalent to the old
        12-pattern regex suite, but 9× faster via Rust. Falls back to html_text_fast
        then to pure-regex fallback (no external deps).
        """
        if not html:
            return ""
        # PRIMARY: nh3 Rust sanitizer + tag stripping (9× faster than pure regex)
        if NH3_AVAILABLE and nh3 is not None:
            try:
                cleaned = nh3.clean(html)
                # nh3 strips dangerous tags (script, style, event handlers) but
                # preserves structural tags — strip all remaining tags for text output
                text = _TAG_STRIP_PATTERN.sub(" ", cleaned)
                return _MULTI_SPACE_PATTERN.sub(" ", text).strip()
            except Exception:  # noqa: BLE001
                pass
        # FALLBACK: html_text_fast canonical helper
        if _CANONICAL_HTML_TEXT_AVAILABLE and html_to_text_fast is not None:
            try:
                return html_to_text_fast(html) or ""
            except Exception:  # noqa: BLE001
                pass
        # ULTIMATE FALLBACK: regex patterns (no external dependencies)
        try:
            text = html
            for pattern, replacement in _NH3_FALLBACK_PATTERNS:
                text = pattern.sub(replacement, text)
            return text.strip()
        except Exception as e:
            logger.warning(f"HTML cleaning failed: {e}")
            return ""

    def _extract_title_fallback(self, html: str) -> str:
        """Extract title using regex (fallback)"""
        try:
            title_match = re.search("<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
            if title_match:
                title = title_match.group(1)
                return re.sub("\\s+", " ", title).strip()
            return ""
        except Exception:
            return ""

    def mine_text(self, text: str, url: str = "") -> MiningResult:
        """
        Mine plain text (minimal processing).

        Args:
            text: Plain text content
            url: Source URL

        Returns:
            MiningResult with cleaned text
        """
        try:
            cleaned = re.sub("\\s+", " ", text).strip()
            return MiningResult(
                content=cleaned, url=url, metadata={"source": url, "method": "plain_text", "char_count": len(cleaned)}
            )
        except Exception as e:
            logger.error(f"Text mining failed: {e}")
            return MiningResult(content="", url=url, success=False, error=str(e))

    def _score_link(self, href: str, base_domain: str, rel_flags: list[str]) -> float:
        """Calculate link score (0-1)."""
        score = 0.5
        link_domain = urlparse(href).netloc.lower()
        if link_domain and link_domain != base_domain:
            score += 0.2
        if href.lower().endswith((".pdf", ".json", ".xml", ".doc", ".docx")):
            score += 0.3
        if "nofollow" in rel_flags or "sponsored" in rel_flags or "ugc" in rel_flags:
            score -= 0.2
        return max(0.0, min(1.0, score))

    # OSINT-03: Maximum HTML input size (5 MB) for selectolax parsing.
    # Prevents OOM on M1 8GB by bounding DOM node allocation.
    _MAX_HTML_INPUT_SIZE: int = 5 * 1024 * 1024

    def _extract_links_selectolax(self, html: str, base_url: str, max_links: int = 50) -> list[dict[str, Any]]:
        """Extract links using selectolax (fast, safe CSS selectors)."""
        if not SELECTOLAX_AVAILABLE:
            return []
        # OSINT-03: Bound input size before parsing to prevent OOM on M1 8GB.
        if len(html) > self._MAX_HTML_INPUT_SIZE:
            html = html[: self._MAX_HTML_INPUT_SIZE]
        try:
            parser = HTMLParser(html)
            links = []
            base_domain = urlparse(base_url).netloc.lower() if base_url else ""
            for node in parser.css("a"):
                if len(links) >= max_links:
                    break
                href = node.attributes.get("href", "").strip()
                if not href or href.startswith(("#", "javascript:", "mailto:")):
                    continue
                if base_url and href.startswith("/"):
                    href = urljoin(base_url, href)
                text = node.text(deep=True).strip()[:120]
                rel = node.attributes.get("rel", "")
                rel_flags = rel.split() if rel else []
                score = self._score_link(href, base_domain, rel_flags)
                links.append(
                    {
                        "url": href,
                        "anchor_text": text,
                        "context_snippet": "",
                        "rel_flags": rel_flags,
                        "score": round(score, 2),
                    }
                )
            return links
        except Exception as e:
            logger.warning(f"selectolax extraction failed: {e}")
            return []

    def extract_jsonld(self, html: str, max_bytes: int = 20000) -> list[dict]:
        """Extract JSON-LD script blocks from HTML."""
        import json

        results = []
        pattern = re.compile(
            "<script[^>]+type=[\"\\']application/ld\\+json[\"\\'][^>]*>(.*?)</script>", re.IGNORECASE | re.DOTALL
        )
        for match in pattern.finditer(html):
            script_content = match.group(1).strip()
            if len(script_content) > max_bytes:
                script_content = script_content[:max_bytes]
            try:
                data = json.loads(script_content)
                results.append(data)
            except json.JSONDecodeError:
                continue
        return results

    def batch_mine(self, html_list: list[tuple[str, str]], include_metadata: bool = False) -> list[MiningResult]:
        """
        Batch mine multiple HTML documents.

        Args:
            html_list: List of (html_content, url) tuples
            include_metadata: Extract metadata

        Returns:
            List of MiningResult objects
        """
        results = []
        for html_content, url in html_list:
            result = self.mine_html(html_content, url, include_metadata)
            results.append(result)
        return results

    def extract_links(self, html_content: str, base_url: str = "", max_links: int = 50) -> list[dict[str, str]]:
        """
        Extract links from HTML with anchor context and scoring - M1 8GB optimized.

        Args:
            html_content: Raw HTML content
            base_url: Base URL for resolving relative links
            max_links: Maximum number of links to extract (hard limit)

        Returns:
            List of dicts with 'url', 'anchor_text', 'context_snippet', 'rel_flags', 'score'
        """
        try:
            if not html_content:
                return []
            # OSINT-03: Bound input size before lxml fallback parsing.
            if len(html_content) > self._MAX_HTML_INPUT_SIZE:
                html_content = html_content[: self._MAX_HTML_INPUT_SIZE]
            links = self._extract_links_selectolax(html_content, base_url, max_links)
            if links:
                return links
            from urllib.parse import urlparse

            links = []
            base_domain = urlparse(base_url).netloc.lower() if base_url else ""
            if LXML_AVAILABLE:
                try:
                    tree = lxml_html.fromstring(html_content)
                    hrefs = tree.xpath("//a/@href")
                    for href in hrefs:
                        if len(links) >= max_links:
                            break
                        if href.startswith(("#", "javascript:", "mailto:", "tel:")):
                            continue
                        if base_url and href.startswith("/"):
                            from urllib.parse import urljoin

                            href = urljoin(base_url, href)
                        try:
                            text_elem = tree.xpath(f'//a[@href="{href}"]/text()')
                            text = " ".join(text_elem).strip()[:120] if text_elem else ""
                        except Exception:
                            text = ""
                        if any(l["url"] == href for l in links):
                            continue
                        try:
                            rel_elem = tree.xpath(f'//a[@href="{href}"]/@rel')
                            rel_flags = rel_elem[0].split() if rel_elem else []
                        except Exception:
                            rel_flags = []
                        score = 0.5
                        link_domain = urlparse(href).netloc.lower()
                        if link_domain and link_domain != base_domain:
                            score += 0.2
                        if href.lower().endswith((".pdf", ".json", ".xml", ".doc", ".docx")):
                            score += 0.3
                        if "nofollow" in rel_flags or "sponsored" in rel_flags or "ugc" in rel_flags:
                            score -= 0.2
                        score = max(0.0, min(1.0, score))
                        links.append(
                            {
                                "url": href,
                                "anchor_text": text,
                                "context_snippet": "",
                                "rel_flags": rel_flags,
                                "score": round(score, 2),
                            }
                        )
                    logger.debug(f"Extracted {len(links)} links via lxml from {base_url}")
                    return links
                except Exception as e:
                    logger.warning(f"lxml parsing failed, falling back to regex: {e}")
            pattern = "<a\\s+([^>]+)>([^<]*)</a>"
            for match in re.finditer(pattern, html_content, re.IGNORECASE):
                if len(links) >= max_links:
                    logger.debug(f"Link limit reached: {max_links}")
                    break
                attrs = match.group(1)
                text = re.sub("\\s+", " ", match.group(2).strip())
                href_match = re.search("href\\s*=\\s*[\"\\']([^\"\\']+)[\"\\']", attrs, re.I)
                if not href_match:
                    continue
                href = href_match.group(1).strip()
                if href.startswith(("#", "javascript:", "mailto:", "tel:")):
                    continue
                if base_url and href.startswith("/"):
                    from urllib.parse import urljoin

                    href = urljoin(base_url, href)
                rel_flags = []
                rel_match = re.search("rel=[\"\\']([^\"\\']+)[\"\\']", attrs, re.I)
                if rel_match:
                    rel_value = rel_match.group(1).lower()
                    rel_flags = rel_value.split()
                start_pos = max(0, match.start() - 200)
                end_pos = min(len(html_content), match.end() + 200)
                context = html_content[start_pos:end_pos]
                context = re.sub("\\s+", " ", context).strip()
                if any(l["url"] == href for l in links):
                    continue
                score = 0.5
                link_domain = urlparse(href).netloc.lower()
                if link_domain and link_domain != base_domain:
                    score += 0.2
                if href.lower().endswith((".pdf", ".json", ".xml", ".doc", ".docx")):
                    score += 0.3
                if "nofollow" in rel_flags or "sponsored" in rel_flags or "ugc" in rel_flags:
                    score -= 0.2
                score = max(0.0, min(1.0, score))
                links.append(
                    {
                        "url": href,
                        "anchor_text": text[:120] if text else "",
                        "context_snippet": context[:200] if context else "",
                        "rel_flags": rel_flags,
                        "score": round(score, 2),
                    }
                )
            logger.debug(f"Extracted {len(links)} links with scoring from {base_url}")
            return links
        except Exception as e:
            logger.warning(f"Link extraction failed: {e}")
            return []


def extract_source_map_url(html: str) -> str | None:
    """
    Find //# sourceMappingURL= in HTML (usually on last lines).
    Returns URL or None.
    """
    tail = html[-5000:]
    match = re.search("//# sourceMappingURL=([^\\s]+)", tail)
    if match:
        url = match.group(1).strip()
        if len(url) > 500:
            domain = url.split("/")[2] if "://" in url else "unknown"
            logger.warning(f"Source map URL from {domain} exceeds 500 chars, truncating")
            url = url[:500]
        return url
    return None


def extract_embedded_json(
    html_content: str,
    url: str = "",
    max_scripts: int = 3,
    max_bytes_per_script: int = 10240,
    max_total_chars: int = 2000,
) -> dict[str, Any]:
    """
    Extract embedded JSON states from HTML (Next.js, React, etc.)

    Extracts:
    - <script id="__NEXT_DATA__" type="application/json">
    - <script type="application/json"> (limited)

    Args:
        html_content: Raw HTML content
        url: Source URL for logging
        max_scripts: Maximum JSON scripts to extract (default: 3)
        max_bytes_per_script: Max bytes per script (default: 10KB)
        max_total_chars: Max total extracted characters (default: 2000)

    Returns:
        Dict with 'embedded_state' containing type, preview, size, extracted_chars
    """
    result = {"embedded_state": None, "extracted_texts": []}
    if not html_content:
        return result
    try:
        import json as json_module

        next_data_pattern = "<script[^>]*id=[\"\\']__NEXT_DATA__[\"\\'][^>]*>(.*?)</script>"
        next_match = re.search(next_data_pattern, html_content, re.DOTALL | re.I)
        if next_match:
            json_str = next_match.group(1).strip()
            if len(json_str) <= max_bytes_per_script:
                try:
                    data = json_module.loads(json_str)
                    texts = _extract_strings_from_json(data, min_len=20, max_len=300)
                    total_chars = sum(len(t) for t in texts)
                    limited_texts = []
                    for t in texts:
                        if total_chars <= max_total_chars:
                            limited_texts.append(t)
                            total_chars += len(t)
                        else:
                            break
                    result["embedded_state"] = {
                        "type": "next_data",
                        "preview": json_str[:500],
                        "size": len(json_str),
                        "extracted_chars": sum(len(t) for t in limited_texts),
                    }
                    result["extracted_texts"] = limited_texts
                    if result is not None and "embedded_state" in result:
                        logger.info(
                            f"[EMBEDDED JSON] url={url} kind=next_data bytes={len(json_str)} extracted_chars={result['embedded_state']['extracted_chars']}"
                        )
                except json_module.JSONDecodeError:  # noqa: BLE001
                    pass
        if len(result["extracted_texts"]) == 0 or len(result["extracted_texts"]) < max_scripts:
            json_pattern = "<script[^>]*type=[\"\\']application/json[\"\\'][^>]*>(.*?)</script>"
            for i, match in enumerate(re.finditer(json_pattern, html_content, re.DOTALL | re.I)):
                if i >= max_scripts:
                    break
                json_str = match.group(1).strip()
                if len(json_str) <= max_bytes_per_script:
                    try:
                        data = json_module.loads(json_str)
                        texts = _extract_strings_from_json(data, min_len=20, max_len=300)
                        total_chars = sum(len(t) for t in texts)
                        limited_texts = []
                        for t in texts:
                            if total_chars <= max_total_chars:
                                limited_texts.append(t)
                                total_chars += len(t)
                            else:
                                break
                        if not result["embedded_state"]:
                            result["embedded_state"] = {
                                "type": "json_script",
                                "preview": json_str[:500],
                                "size": len(json_str),
                                "extracted_chars": sum(len(t) for t in limited_texts),
                            }
                            result["extracted_texts"] = limited_texts
                        logger.info(
                            f"[EMBEDDED JSON] url={url} kind=json_script_{i} bytes={len(json_str)} extracted_chars={sum(len(t) for t in limited_texts)}"
                        )
                    except json_module.JSONDecodeError:  # noqa: BLE001
                        pass
    except Exception as e:
        logger.debug(f"Embedded JSON extraction failed: {e}")
    return result


def _extract_strings_from_json(
    obj: Any, min_len: int = 20, max_len: int = 300, max_depth: int = 10, current_depth: int = 0
) -> list[str]:
    """
    Recursively extract string values from JSON that look like content (20-300 chars).

    Args:
        obj: JSON object (dict, list, or primitive)
        min_len: Minimum string length
        max_len: Maximum string length
        max_depth: Maximum recursion depth
        current_depth: Current recursion depth

    Returns:
        List of extracted strings
    """
    texts = []
    if current_depth > max_depth:
        return texts
    if isinstance(obj, str):
        if min_len <= len(obj) <= max_len:
            if not obj.startswith(("http://", "https://", "/", "./", "../")):
                if not re.match("^[\\d\\.\\-\\+]+$", obj):
                    texts.append(obj)
    elif isinstance(obj, dict):
        skip_keys = {"props", "pageProps", "initialState", "__typename", "id", "name", "type", "key"}
        for key, value in obj.items():
            if key.lower() not in skip_keys:
                texts.extend(_extract_strings_from_json(value, min_len, max_len, max_depth, current_depth + 1))
    elif isinstance(obj, (list, tuple)):
        for item in obj:
            if len(texts) < 50:
                texts.extend(_extract_strings_from_json(item, min_len, max_len, max_depth, current_depth + 1))
    return texts


def create_rust_miner(prefer_rust: bool = True) -> RustMiner:
    """
    Factory function to create a RustMiner.

    Args:
        prefer_rust: Prefer Rust-based libraries

    Returns:
        RustMiner instance
    """
    return RustMiner(prefer_rust=prefer_rust)


class FeedDiscoveryResult(Struct, frozen=True):
    """Result of feed discovery."""

    feed_urls: list[str]
    source_url: str
    discovery_method: str


class FeedDiscoverer:
    """Discover RSS/Atom feeds from HTML content."""

    __slots__ = ("max_heuristic_feeds",)

    def __init__(self, max_heuristic_feeds: int = 10) -> None:
        self.max_heuristic_feeds = max_heuristic_feeds

    def discover_feeds(self, html_content: str, base_url: str = "") -> FeedDiscoveryResult:
        """
        Discover RSS/Atom feeds in HTML content.

        Args:
            html_content: Raw HTML content
            base_url: Base URL for resolving relative links

        Returns:
            FeedDiscoveryResult with discovered feed URLs
        """
        feed_urls = []
        link_feeds = self._extract_from_link_tags(html_content, base_url)
        feed_urls.extend(link_feeds)
        if len(feed_urls) < 3:
            heuristic_feeds = self._heuristic_discovery(html_content, base_url)
            for feed in heuristic_feeds:
                if feed not in feed_urls:
                    feed_urls.append(feed)
                    if len(feed_urls) >= self.max_heuristic_feeds:
                        break
        discovery_method = "link_tag" if link_feeds else "heuristic"
        if link_feeds and len(feed_urls) > len(link_feeds):
            discovery_method = "mixed"
        return FeedDiscoveryResult(
            feed_urls=feed_urls[: self.max_heuristic_feeds], source_url=base_url, discovery_method=discovery_method
        )

    def _extract_from_link_tags(self, html_content: str, base_url: str) -> list[str]:
        """Extract feed URLs from <link rel="alternate"> tags."""
        feeds = []
        pattern = "<link[^>]+rel=[\"\\']alternate[\"\\'][^>]+type=[\"\\'](application/rss\\+xml|application/atom\\+xml|application/feed\\+json)[\"\\'][^>]+href=[\"\\']([^\"\\']+)[\"\\']"
        for match in re.finditer(pattern, html_content, re.IGNORECASE):
            href = match.group(2).strip()
            resolved = self._resolve_url(href, base_url)
            if resolved and resolved not in feeds:
                feeds.append(resolved)
        pattern2 = "<link[^>]+type=[\"\\'](application/rss\\+xml|application/atom\\+xml|application/feed\\+json)[\"\\'][^>]+rel=[\"\\']alternate[\"\\'][^>]+href=[\"\\']([^\"\\']+)[\"\\']"
        for match in re.finditer(pattern2, html_content, re.IGNORECASE):
            href = match.group(2).strip()
            resolved = self._resolve_url(href, base_url)
            if resolved and resolved not in feeds:
                feeds.append(resolved)
        return feeds

    def _heuristic_discovery(self, html_content: str, base_url: str) -> list[str]:
        """Heuristic feed discovery based on common paths."""
        from urllib.parse import urljoin, urlparse

        if not base_url:
            return []
        parsed = urlparse(base_url)
        base = f"{parsed.scheme}://{parsed.netloc}"
        common_paths = [
            "/feed",
            "/feed.xml",
            "/feed/",
            "/rss",
            "/rss.xml",
            "/rss/",
            "/atom",
            "/atom.xml",
            "/atom/",
            "/index.xml",
            "/posts.xml",
            "/blog/feed",
            "/blog/rss",
            "/blog/atom",
            "/feeds/posts/default",
            "/feeds/posts/default?alt=rss",
            "/?feed=rss2",
            "/?feed=atom",
            "/wp-feed.php",
            "/.rss",
            "/.atom",
            "/jsonfeed.json",
            "/feed.json",
        ]
        return [urljoin(base, path) for path in common_paths]

    def _resolve_url(self, href: str, base_url: str) -> str:
        """Resolve relative URL to absolute."""
        from urllib.parse import urljoin

        if not href:
            return ""
        if href.startswith(("http://", "https://")):
            return href
        if base_url:
            return urljoin(base_url, href)
        return ""


class ExtractedMetadata(Struct, frozen=True):
    """Metadata extracted from non-HTML documents."""

    content_type: str
    file_size: int
    title: str | None = None
    author: str | None = None
    creation_date: str | None = None
    modification_date: str | None = None
    page_count: int | None = None
    keywords: list[str] = field(default_factory=list)
    entities: list[dict[str, Any]] = field(default_factory=list)
    gps_coords: tuple[float, float] | None = None
    timeline_events: list[dict[str, Any]] = field(default_factory=list)
    extracted_text_preview: str = ""


class MetadataExtractor:
    """Extract metadata from non-HTML documents (PDF, images, etc.) - M1 8GB."""

    __slots__ = ("_exifread_available", "_pillow_available", "_pymupdf_available")

    def __init__(self) -> None:
        self._pymupdf_available = None
        self._exifread_available = None
        self._pillow_available = None

    def _check_pymupdf(self) -> bool:
        if self._pymupdf_available is None:
            try:
                import fitz

                self._pymupdf_available = True
            except ImportError:
                self._pymupdf_available = False
        return self._pymupdf_available

    def _check_exifread(self) -> bool:
        if self._exifread_available is None:
            try:
                import exifread

                self._exifread_available = True
            except ImportError:
                self._exifread_available = False
        return self._exifread_available

    def _check_pillow(self) -> bool:
        if self._pillow_available is None:
            try:
                from PIL import Image

                self._pillow_available = True
            except ImportError:
                self._pillow_available = False
        return self._pillow_available

    async def extract(self, content_bytes: bytes, content_type: str) -> ExtractedMetadata:
        """Extract metadata based on content-type."""
        metadata = ExtractedMetadata(content_type=content_type, file_size=len(content_bytes))
        if "application/pdf" in content_type:
            metadata = await self._extract_pdf(content_bytes, metadata)
        elif content_type.startswith("image/"):
            metadata = await self._extract_image(content_bytes, metadata)
        elif content_type in [
            "application/msword",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ]:
            metadata = await self._extract_docx(content_bytes, metadata)
        elif content_type in [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ]:
            metadata = await self._extract_xlsx(content_bytes, metadata)
        elif content_type in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ]:
            metadata = await self._extract_pptx(content_bytes, metadata)
        return metadata

    async def _extract_pdf(self, content_bytes: bytes, metadata: ExtractedMetadata) -> ExtractedMetadata:
        """Extract metadata from PDF.

        Strategy:
        1. Try Rust lopdf via hledac_rust_extensions (fast, pure Rust)
        2. Fallback to PyMuPDF (fitz) if Rust is unavailable
        """
        # Try Rust PDF extraction first (feature-gated, ~10× faster than Python pypdf)
        _rust_pdf_available = False
        # R6: Centralized Rust access via core.rust_backend
        from hledac.universal._core.rust_backend import rust

        if rust.is_available and rust.raw.pdf is not None:
            _rust_pdf_available = True

        if _rust_pdf_available:
            try:
                import asyncio

                def _extract_rust():
                    # Single-pass: extract text + IOCs together (avoids 2× PDF parsing)
                    text, ioc_list = rust.raw.pdf.extract_text_and_iocs_from_bytes(content_bytes)
                    return text, len(ioc_list)

                asyncio.get_running_loop()
                text, ioc_count = await asyncio.to_thread(_extract_rust)
                text_preview = text[:2000] if text else ""
                logger.debug(f"[METADATA] Rust PDF extracted: {ioc_count} IOCs, preview length: {len(text_preview)}")
                return ExtractedMetadata(
                    content_type=metadata.content_type,
                    file_size=metadata.file_size,
                    title=metadata.title,
                    author=metadata.author,
                    creation_date=metadata.creation_date,
                    modification_date=metadata.modification_date,
                    page_count=metadata.page_count,
                    keywords=metadata.keywords,
                    entities=metadata.entities,
                    gps_coords=metadata.gps_coords,
                    timeline_events=metadata.timeline_events,
                    extracted_text_preview=text_preview,
                )
            except Exception as e:
                logger.debug(f"[METADATA] Rust PDF extraction failed, falling back to PyMuPDF: {e}")
                # Fall through to PyMuPDF fallback

        # PyMuPDF fallback
        if not self._check_pymupdf():
            return metadata
        try:
            import asyncio

            import fitz

            def _extract():
                doc = fitz.open(stream=content_bytes, filetype="pdf")
                meta = doc.metadata
                result = {
                    "title": meta.get("title"),
                    "author": meta.get("author"),
                    "creation_date": meta.get("creationDate"),
                    "modification_date": meta.get("modDate"),
                    "page_count": len(doc),
                    "keywords": meta.get("keywords", "").split(",") if meta.get("keywords") else [],
                    "text_preview": "",
                }
                if len(doc) > 0:
                    text_parts = [doc[i].get_text() for i in range(len(doc))]
                    text = "".join(text_parts)
                    result["text_preview"] = text[:2000]
                doc.close()
                return result

            asyncio.get_running_loop()
            result = await asyncio.to_thread(_extract)
            logger.debug(f"[METADATA] Extracted PDF: {result['title'] or 'no title'}, {result['page_count']} pages")
            return ExtractedMetadata(
                content_type=metadata.content_type,
                file_size=metadata.file_size,
                title=result["title"],
                author=result["author"],
                creation_date=result["creation_date"],
                modification_date=result["modification_date"],
                page_count=result["page_count"],
                keywords=result["keywords"],
                entities=metadata.entities,
                gps_coords=metadata.gps_coords,
                timeline_events=metadata.timeline_events,
                extracted_text_preview=result["text_preview"],
            )
        except Exception as e:
            logger.warning(f"PDF metadata extraction failed: {e}")
        return metadata

    async def _extract_image(self, content_bytes: bytes, metadata: ExtractedMetadata) -> ExtractedMetadata:
        """Extract metadata from images (EXIF)."""
        if not self._check_exifread():
            return metadata
        try:
            import asyncio
            from io import BytesIO

            import exifread

            def _extract():
                result = {"gps_coords": None, "creation_date": None, "camera_model": None}
                tags = exifread.process_file(BytesIO(content_bytes), details=False)
                if "GPS GPSLatitude" in tags and "GPS GPSLongitude" in tags:
                    try:
                        lat = self._convert_gps(tags["GPS GPSLatitude"])
                        lon = self._convert_gps(tags["GPS GPSLongitude"])
                        if "GPS GPSLatitudeRef" in tags and str(tags["GPS GPSLatitudeRef"]) == "S":
                            lat = -lat
                        if "GPS GPSLongitudeRef" in tags and str(tags["GPS GPSLongitudeRef"]) == "W":
                            lon = -lon
                        result["gps_coords"] = (lat, lon)
                    except Exception:  # noqa: BLE001
                        pass
                if "EXIF DateTimeOriginal" in tags:
                    result["creation_date"] = str(tags["EXIF DateTimeOriginal"])
                return result

            asyncio.get_running_loop()
            result = await asyncio.to_thread(_extract)
            metadata.gps_coords = result["gps_coords"]
            if result["creation_date"]:
                metadata.creation_date = result["creation_date"]
        except Exception as e:
            logger.warning(f"Image metadata extraction failed: {e}")
        return metadata

    def _convert_gps(self, gps_tag) -> float:
        """Convert EXIF GPS coordinates to decimal."""
        values = gps_tag.values
        d = float(values[0].num) / float(values[0].den)
        m = float(values[1].num) / float(values[1].den)
        s = float(values[2].num) / float(values[2].den)
        return d + m / 60.0 + s / 3600.0

    async def _extract_docx(self, content_bytes: bytes, metadata: ExtractedMetadata) -> ExtractedMetadata:
        """Extract metadata from DOCX/XLSX/PPTX.

        Strategy:
        1. Try Rust office extraction via hledac_rust_extensions (fast, pure Rust)
        2. Fallback to python-docx/openpyxl if Rust is unavailable
        """
        # Try Rust office extraction first (feature-gated, ~5-10× faster than Python)
        _rust_office_available = False
        from hledac.universal._core.rust_backend import rust

        if rust.is_available and rust.raw.office is not None:
            _rust_office_available = True

        if _rust_office_available:
            try:
                import asyncio

                def _extract_rust():
                    # Returns (text, ioc_count) or raises
                    text = rust.raw.office.extract_text_from_bytes(content_bytes, "docx")
                    ioc_list = rust.raw.office.extract_iocs_from_bytes(content_bytes, "docx")
                    return text, len(ioc_list)

                asyncio.get_running_loop()
                text, ioc_count = await asyncio.to_thread(_extract_rust)
                text_preview = text[:2000] if text else ""
                logger.debug(f"[METADATA] Rust DOCX extracted: {ioc_count} IOCs, preview length: {len(text_preview)}")
                return msgspec.replace(metadata, extracted_text_preview=text_preview)
            except Exception as e:
                logger.debug(f"[METADATA] Rust DOCX extraction failed, falling back to python-docx: {e}")
                # Fall through to python-docx fallback

        # python-docx fallback
        try:
            import asyncio

            import docx

            def _extract():
                doc = docx.Document(io.BytesIO(content_bytes))
                text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text)
                return " ".join(text_parts)

            asyncio.get_running_loop()
            text = await asyncio.to_thread(_extract)
            text_preview = text[:2000] if text else ""
            return msgspec.replace(metadata, extracted_text_preview=text_preview)
        except Exception as e:
            logger.debug(f"[METADATA] DOCX extraction failed: {e}")
        return metadata

    async def _extract_xlsx(self, content_bytes: bytes, metadata: ExtractedMetadata) -> ExtractedMetadata:
        """Extract metadata from XLSX.

        Strategy:
        1. Try Rust office extraction via hledac_rust_extensions (fast, pure Rust)
        2. Fallback to openpyxl if Rust is unavailable
        """
        _rust_office_available = False
        from hledac.universal._core.rust_backend import rust

        if rust.is_available and rust.raw.office is not None:
            _rust_office_available = True

        if _rust_office_available:
            try:
                import asyncio

                def _extract_rust():
                    text = rust.raw.office.extract_text_from_bytes(content_bytes, "xlsx")
                    ioc_list = rust.raw.office.extract_iocs_from_bytes(content_bytes, "xlsx")
                    return text, len(ioc_list)

                asyncio.get_running_loop()
                text, ioc_count = await asyncio.to_thread(_extract_rust)
                text_preview = text[:2000] if text else ""
                logger.debug(f"[METADATA] Rust XLSX extracted: {ioc_count} IOCs, preview length: {len(text_preview)}")
                return msgspec.replace(metadata, extracted_text_preview=text_preview)
            except Exception as e:
                logger.debug(f"[METADATA] Rust XLSX extraction failed, falling back to openpyxl: {e}")

        # openpyxl fallback
        try:
            import asyncio

            import openpyxl

            def _extract():
                wb = openpyxl.load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
                text_parts = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(max_row=100):
                        for cell in row:
                            if cell.value and isinstance(cell.value, str):
                                text_parts.append(cell.value)
                wb.close()
                return " ".join(text_parts)

            asyncio.get_running_loop()
            text = await asyncio.to_thread(_extract)
            text_preview = text[:2000] if text else ""
            return msgspec.replace(metadata, extracted_text_preview=text_preview)
        except Exception as e:
            logger.debug(f"[METADATA] XLSX extraction failed: {e}")
            return metadata

    async def _extract_pptx(self, content_bytes: bytes, metadata: ExtractedMetadata) -> ExtractedMetadata:
        """Extract metadata from PPTX.

        Strategy:
        1. Try Rust office extraction via hledac_rust_extensions (fast, pure Rust)
        2. Fallback to python-pptx if Rust is unavailable
        """
        _rust_office_available = False
        from hledac.universal._core.rust_backend import rust

        if rust.is_available and rust.raw.office is not None:
            _rust_office_available = True

        if _rust_office_available:
            try:
                import asyncio

                def _extract_rust():
                    text = rust.raw.office.extract_text_from_bytes(content_bytes, "pptx")
                    ioc_list = rust.raw.office.extract_iocs_from_bytes(content_bytes, "pptx")
                    return text, len(ioc_list)

                asyncio.get_running_loop()
                text, ioc_count = await asyncio.to_thread(_extract_rust)
                text_preview = text[:2000] if text else ""
                logger.debug(f"[METADATA] Rust PPTX extracted: {ioc_count} IOCs, preview length: {len(text_preview)}")
                return msgspec.replace(metadata, extracted_text_preview=text_preview)
            except Exception as e:
                logger.debug(f"[METADATA] Rust PPTX extraction failed, falling back to python-pptx: {e}")

        # python-pptx fallback
        try:
            import asyncio

            from pptx import Presentation

            def _extract():
                prs = Presentation(io.BytesIO(content_bytes))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                return " ".join(text_parts)

            asyncio.get_running_loop()
            text = await asyncio.to_thread(_extract)
            text_preview = text[:2000] if text else ""
            return msgspec.replace(metadata, extracted_text_preview=text_preview)
        except Exception as e:
            logger.debug(f"[METADATA] PPTX extraction failed: {e}")
            return metadata


import ast
import hashlib
import os
import sys
import time
from collections import OrderedDict
from concurrent.futures import ThreadPoolExecutor


def build_structure_map(root_dir: str, *, limits: dict, state: dict) -> dict:
    """
    Build structure map: scan Python project, extract imports, build dependency graph.

    Args:
        root_dir: Root directory to scan
        limits: Resource limits (max_files, max_bytes_total, time_budget_ms, etc.)
        state: Persistent state (file_cache LRU, prev_edges)

    Returns:
        Dict with fingerprint, files, edges, meta
    """
    state = state or {}
    start_time = time.monotonic()
    max_files = limits.get("max_files", 2500)
    max_bytes_total = limits.get("max_bytes_total", 8000000)
    max_parse_bytes = limits.get("max_parse_bytes_per_file", 65536)
    time_budget_ms = limits.get("time_budget_ms", 1200)
    prefix_hash_bytes = limits.get("prefix_hash_bytes", 4096)
    incremental = limits.get("incremental", True)
    parallel_threshold = limits.get("parallel_scan_threshold", 5000)
    max_workers = limits.get("max_workers", min(4, os.cpu_count() or 4))
    errors: list[str] = []
    truncated = False
    truncation_reason: str | None = None
    total_bytes = 0
    seen_inodes: set[tuple[int, int]] = set()
    candidates: list[tuple[str, os.DirEntry]] = []

    def _scan_recursive(entry: os.DirEntry, depth: int = 0) -> None:
        nonlocal total_bytes, truncated, truncation_reason
        if len(candidates) >= max_files or total_bytes >= max_bytes_total:
            if not truncated:
                truncated = True
                truncation_reason = "file_budget" if len(candidates) >= max_files else "size_budget"
            return
        if time.monotonic() - start_time > time_budget_ms / 1000:
            truncated = True
            truncation_reason = "time_budget"
            return
        try:
            if entry.is_dir(follow_symlinks=False):
                if entry.name.startswith(".") or entry.name in ("__pycache__", "node_modules", "venv", ".venv"):
                    return
                try:
                    stat = entry.stat(follow_symlinks=False)
                    inode_key = (stat.st_dev, stat.st_ino)
                    if inode_key in seen_inodes:
                        return
                    seen_inodes.add(inode_key)
                except OSError:
                    return
                with os.scandir(entry.path) as it:
                    for sub in it:
                        _scan_recursive(sub, depth + 1)
            elif entry.is_file(follow_symlinks=False):
                if not entry.name.endswith(".py"):
                    return
                try:
                    stat = entry.stat(follow_symlinks=False)
                    inode_key = (stat.st_dev, stat.st_ino)
                    if inode_key in seen_inodes:
                        return
                    seen_inodes.add(inode_key)
                    if stat.st_size == 0:
                        return
                    candidates.append((entry.path, entry, stat))
                    total_bytes += stat.st_size
                except OSError:  # noqa: BLE001
                    pass
        except PermissionError:  # noqa: BLE001
            pass

    try:
        with os.scandir(root_dir) as it:
            for entry in it:
                _scan_recursive(entry)
    except PermissionError:
        errors.append(f"Permission denied: {root_dir}")
    use_parallel = len(candidates) > parallel_threshold
    files_data: list[dict[str, Any]] = []
    file_cache = state.get("file_cache", OrderedDict())

    def _process_file(
        path: str, entry: os.DirEntry, stat_result: os.stat_result | None = None
    ) -> dict[str, Any] | None:
        nonlocal errors
        try:
            stat = stat_result if stat_result is not None else entry.stat(follow_symlinks=False)
            mtime_ns = stat.st_mtime_ns
            size = stat.st_size
            prefix_bytes = _read_prefix_bytes(path, prefix_hash_bytes, errors, stat_result=stat)
            prefix_hash = _hash_bytes(prefix_bytes)
            text = ""
            parse_mode = "ast"
            if prefix_bytes:
                try:
                    text = prefix_bytes.decode("utf-8", errors="replace")
                except Exception:  # noqa: BLE001
                    pass
            imports: list[str] = []
            if text:
                try:
                    tree = ast.parse(text, type_comments=False)
                    imports = _extract_imports_ast(tree)
                    parse_mode = "ast"
                except SyntaxError:
                    imports = _extract_imports_regex(text)
                    parse_mode = "regex"
            rel_path = os.path.relpath(path, root_dir)
            module = _path_to_module(rel_path)
            cache_key = rel_path
            cached = file_cache.get(cache_key, {})
            cached_hash = cached.get("prefix_hash", "")
            changed = prefix_hash != cached_hash
            return {
                "rel_path": rel_path,
                "module": module,
                "mtime_ns": mtime_ns,
                "size": size,
                "prefix_hash": prefix_hash,
                "imports": imports,
                "parse_mode": parse_mode,
                "changed": changed,
            }
        except Exception as e:
            errors.append(f"Error processing {path}: {e}")
            return None

    if use_parallel:
        remaining_ms = time_budget_ms - (time.monotonic() - start_time) * 1000
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            try:
                for result in executor.map(
                    _process_file,
                    [p for p, _, _ in candidates],
                    [e for _, e, _ in candidates],
                    [s for _, _, s in candidates],
                    buffersize=8,
                    timeout=remaining_ms / 1000,
                ):
                    if result:
                        files_data.append(result)
            except TimeoutError:
                truncated = True
                truncation_reason = "time_budget"
            except Exception:  # noqa: BLE001
                pass
    else:
        for path, entry, stat_result in candidates:
            if truncated:
                break
            result = _process_file(path, entry, stat_result)
            if result:
                files_data.append(result)
    changed_files = [f for f in files_data if f.get("changed", False)]
    changed_modules = sorted({f["module"] for f in changed_files if f["module"]})
    for f in files_data:
        rel_path = f["rel_path"]
        file_cache[rel_path] = {
            "imports": f["imports"],
            "prefix_hash": f["prefix_hash"],
            "mtime_ns": f["mtime_ns"],
            "size": f["size"],
            "module": f["module"],
            "hot_score": file_cache.get(rel_path, {}).get("hot_score", 1.0),
            "last_access_ts": time.time(),
            "parse_mode": f["parse_mode"],
        }
    while len(file_cache) > 512:
        file_cache.popitem(last=False)
    module_set = {f["module"] for f in files_data if f["module"]}
    prev_edges = state.get("prev_edges", [])
    edges: list[dict[str, Any]] = []
    if incremental and prev_edges:
        stable_edges = [e for e in prev_edges if e.get("src") not in changed_modules]
        edges.extend(stable_edges)
        for f in files_data:
            if f["module"] in changed_modules:
                for imp in f["imports"]:
                    external = not (imp in module_set or imp.startswith("hledac."))
                    edges.append({"src": f["module"], "dst": imp, "external": external})
    else:
        for f in files_data:
            for imp in f["imports"]:
                external = not (imp in module_set or imp.startswith("hledac."))
                edges.append({"src": f["module"], "dst": imp, "external": external})
    edges.sort(key=lambda e: (e.get("src", ""), e.get("dst", "")))
    elapsed_ms = int((time.monotonic() - start_time) * 1000)
    total_files = len(files_data)
    churn_ratio = len(changed_files) / total_files if total_files > 0 else 0.0
    sorted_files = sorted(files_data, key=itemgetter("rel_path"))
    fingerprint_input = {
        "files": [(f["rel_path"], f["prefix_hash"], f.get("mtime_ns", 0)) for f in sorted_files],
        "edges": [(e["src"], e["dst"]) for e in edges],
        "limits_used": {
            "max_files": max_files,
            "max_bytes_total": max_bytes_total,
            "max_parse_bytes_per_file": max_parse_bytes,
            "time_budget_ms": time_budget_ms,
            "prefix_hash_bytes": prefix_hash_bytes,
            "incremental": incremental,
            "parallel_scan_threshold": parallel_threshold,
            "max_workers": max_workers,
        },
        "version": "1.0",
    }
    fingerprint = _compute_fingerprint(fingerprint_input)
    return {
        "fingerprint": fingerprint,
        "files": files_data,
        "edges": edges,
        "meta": {
            "version": "1.0",
            "limits_used": fingerprint_input["limits_used"],
            "elapsed_ms": elapsed_ms,
            "total_files": total_files,
            "changed_files": len(changed_files),
            "changed_modules": changed_modules,
            "churn_ratio": churn_ratio,
            "truncated": truncated,
            "truncation_reason": truncation_reason,
            "errors": errors,
        },
    }


def _read_prefix_bytes(path: str, n: int, errors: list[str], *, stat_result: os.stat_result | None = None) -> bytes:
    """Read first n bytes with fail-safe."""
    try:
        size = stat_result.st_size if stat_result else 0
        if size == 0:
            return b""
        read_size = min(n, size)
        with open(path, "rb") as f:
            data = f.read(read_size)
            return data
    except PermissionError:
        errors.append(f"Permission denied: {path}")
        return b""
    except (ValueError, OSError) as e:
        errors.append(f"Error reading {path}: {e}")
        return b""


_RUST_XXHASH_AVAILABLE = False
_content_hash_64_rust: Callable[[bytes], int] | None = None
# R6: Centralized Rust access via core.rust_backend
from hledac.universal._core.rust_backend import rust

_content_hash_64_rust = rust.raw.content_hash_64
_RUST_XXHASH_AVAILABLE = _content_hash_64_rust is not None


def _hash_bytes(data: bytes) -> str:
    """Hash bytes using Rust xxhash3-64 (SIMD NEON), Python xxhash, or sha256."""
    if not data:
        return ""
    if _RUST_XXHASH_AVAILABLE and _content_hash_64_rust:
        try:
            return f"{_content_hash_64_rust(data):016x}"
        except Exception:  # noqa: BLE001
            pass
    try:
        import xxhash

        return xxhash.xxh3_64(data).hexdigest()
    except ImportError:  # noqa: BLE001
        pass
    return hashlib.sha256(data).hexdigest()[:16]


def _extract_imports_ast(tree: ast.AST) -> list[str]:
    """Extract imports using AST."""
    imports: list[str] = []
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                imports.append(sys.intern(alias.name))
        elif isinstance(node, ast.ImportFrom):
            module = node.module or ""
            level = node.level
            if level > 0:
                resolved = _resolve_relative_import("", level, module)
                imports.append(sys.intern(resolved))
            else:
                for alias in node.names:
                    full_name = f"{module}.{alias.name}" if module else alias.name
                    imports.append(sys.intern(full_name))
    return imports


def _extract_imports_regex(text: str) -> list[str]:
    """Fallback: extract imports using regex."""
    imports: list[str] = []
    for match in re.finditer("^import\\s+(\\S+)", text, re.MULTILINE):
        imports.append(sys.intern(match.group(1)))
    for match in re.finditer("^from\\s+(\\S+)\\s+import", text, re.MULTILINE):
        module = match.group(1)
        if module.startswith("."):
            resolved = _resolve_relative_import("", module.count("."), module.lstrip("."))
            imports.append(sys.intern(resolved))
        else:
            imports.append(sys.intern(module))
    return imports


def _resolve_relative_import(package_name: str, level: int, module: str) -> str:
    """Resolve relative import to absolute."""
    if level == 0:
        return module
    parts = package_name.split(".")
    if level > len(parts):
        return module
    base_parts = parts[: len(parts) - (level - 1)] if level > 1 else parts
    base = ".".join(base_parts)
    return f"{base}.{module}" if module else base


def _path_to_module(rel_path: str) -> str:
    """Convert file path to module name."""
    if rel_path.endswith("/__init__.py"):
        rel_path = rel_path[:-12]
    elif rel_path.endswith(".py"):
        rel_path = rel_path[:-3]
    parts = rel_path.split("/")
    return ".".join(parts)


def _compute_fingerprint(data: dict) -> str:
    """Compute stable fingerprint from canonical JSON."""
    import json

    canonical = json.dumps(data, ensure_ascii=True, sort_keys=True)
    return hashlib.sha256(canonical.encode()).hexdigest()[:16]
